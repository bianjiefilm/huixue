#!/usr/bin/env python3
"""
课程PPT生成工具
使用python-pptx生成课程PPT
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# 主题颜色
THEME_COLORS = {
    "primary": RGBColor(24, 144, 255),      # 蓝色
    "secondary": RGBColor(82, 196, 26),     # 绿色
    "accent": RGBColor(250, 173, 20),       # 橙色
    "dark": RGBColor(38, 38, 38),           # 深灰
    "light": RGBColor(250, 250, 250),       # 浅灰
    "white": RGBColor(255, 255, 255),
}


def add_title_slide(prs, title, subtitle=""):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    # 背景渐变色块
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME_COLORS["primary"]
    shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = THEME_COLORS["white"]
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(9), Inches(0.8)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = THEME_COLORS["light"]
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, note=""):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 顶部色块
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = THEME_COLORS["primary"]
    header.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = THEME_COLORS["white"]

    # 内容要点
    content_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.6), Inches(8.6), Inches(5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(22)
        p.font.color.rgb = THEME_COLORS["dark"]
        p.space_before = Pt(12)
        p.space_after = Pt(8)

    # 添加备注（用于语音旁白）
    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note

    return slide


def add_summary_slide(prs, title, points):
    """添加总结页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 背景
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME_COLORS["secondary"]
    shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.8), Inches(9), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = THEME_COLORS["white"]
    p.alignment = PP_ALIGN.CENTER

    # 要点
    content_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.2), Inches(8), Inches(4.5)
    )
    tf = content_box.text_frame

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✓ {point}"
        p.font.size = Pt(24)
        p.font.color.rgb = THEME_COLORS["white"]
        p.space_before = Pt(16)

    return slide


def create_chapter_ppt(chapter_data: dict, output_path: str):
    """
    创建章节PPT

    Args:
        chapter_data: 章节数据，包含title, slides等
        output_path: 输出路径
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 标题页
    add_title_slide(
        prs,
        chapter_data["title"],
        chapter_data.get("subtitle", "神经网络与深度学习")
    )

    # 内容页
    for slide_data in chapter_data.get("slides", []):
        add_content_slide(
            prs,
            slide_data["title"],
            slide_data["bullets"],
            slide_data.get("note", "")
        )

    # 总结页
    if "summary" in chapter_data:
        add_summary_slide(prs, "本章小结", chapter_data["summary"])

    # 保存
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"PPT已保存: {output_path}")


# 神经网络与深度学习课程章节数据
CHAPTERS_DATA = [
    {
        "id": "01",
        "title": "人工智能起源与发展",
        "subtitle": "第一章 · 神经网络与深度学习",
        "slides": [
            {
                "title": "什么是人工智能？",
                "bullets": [
                    "人工智能(AI)是让机器模拟人类智能的技术",
                    "包括学习、推理、感知、理解语言等能力",
                    "目标是创造能够自主思考和决策的智能系统",
                    "应用领域：图像识别、语音处理、自然语言处理等"
                ],
                "note": "人工智能是当今最热门的技术领域之一。它的核心目标是让计算机能够像人类一样思考、学习和做出决策。从智能手机的语音助手到自动驾驶汽车，人工智能已经深入我们生活的方方面面。"
            },
            {
                "title": "人工智能发展历程",
                "bullets": [
                    "1956年：达特茅斯会议，AI概念首次提出",
                    "1960-1970年代：专家系统兴起",
                    "1980-1990年代：机器学习方法发展",
                    "2006年至今：深度学习革命，GPU加速训练"
                ],
                "note": "人工智能的发展并非一帆风顺。从1956年达特茅斯会议首次提出人工智能概念，到经历多次寒冬和复兴，直到2006年深度学习的突破，AI终于迎来了真正的黄金时代。"
            },
            {
                "title": "人工智能的三个层次",
                "bullets": [
                    "弱人工智能：专注于特定任务（如下棋、图像识别）",
                    "强人工智能：具有人类级别的通用智能",
                    "超级人工智能：超越人类智能的假设阶段",
                    "目前我们主要处于弱人工智能阶段"
                ],
                "note": "人工智能可以分为三个层次。目前我们所使用的AI都属于弱人工智能，它们在特定任务上可能超越人类，但无法进行通用的智能活动。强人工智能和超级人工智能仍是未来的研究方向。"
            }
        ],
        "summary": [
            "理解人工智能的基本概念和发展历史",
            "掌握AI的三个发展层次",
            "了解深度学习在AI发展中的重要地位"
        ]
    },
    {
        "id": "02",
        "title": "TensorFlow环境安装",
        "subtitle": "第二章 · 神经网络与深度学习",
        "slides": [
            {
                "title": "TensorFlow简介",
                "bullets": [
                    "Google开发的开源深度学习框架",
                    "支持CPU、GPU、TPU等多种硬件加速",
                    "TensorFlow 2.0采用即时执行模式，更易上手",
                    "广泛应用于学术研究和工业生产"
                ],
                "note": "TensorFlow是目前最流行的深度学习框架之一，由Google开发并开源。TensorFlow 2.0版本大幅简化了使用方式，采用即时执行模式，让深度学习编程变得更加直观和高效。"
            },
            {
                "title": "环境准备",
                "bullets": [
                    "Python 3.8或更高版本",
                    "推荐使用Anaconda管理环境",
                    "GPU版本需要CUDA和cuDNN支持",
                    "建议使用虚拟环境隔离项目依赖"
                ],
                "note": "在安装TensorFlow之前，我们需要准备好Python环境。推荐使用Anaconda来管理Python环境，它可以方便地创建和管理多个虚拟环境，避免不同项目之间的依赖冲突。"
            },
            {
                "title": "安装步骤",
                "bullets": [
                    "创建虚拟环境：conda create -n tf python=3.9",
                    "激活环境：conda activate tf",
                    "安装TensorFlow：pip install tensorflow",
                    "验证安装：import tensorflow as tf"
                ],
                "note": "安装TensorFlow非常简单。首先创建一个专门的虚拟环境，然后使用pip命令安装即可。安装完成后，我们可以通过导入tensorflow模块来验证安装是否成功。"
            },
            {
                "title": "验证与测试",
                "bullets": [
                    "检查版本：tf.__version__",
                    "检查GPU支持：tf.config.list_physical_devices('GPU')",
                    "运行简单计算：tf.constant([1, 2, 3])",
                    "完成第一个TensorFlow程序"
                ],
                "note": "安装完成后，我们需要验证TensorFlow是否正常工作。可以检查版本号、GPU支持情况，并运行一些简单的张量计算来确保一切正常。"
            }
        ],
        "summary": [
            "了解TensorFlow框架的特点和优势",
            "掌握TensorFlow 2.0环境搭建方法",
            "能够验证环境配置是否正确"
        ]
    },
    {
        "id": "03",
        "title": "Python语言基础",
        "subtitle": "第三章 · 神经网络与深度学习",
        "slides": [
            {
                "title": "Python基础语法",
                "bullets": [
                    "变量定义与数据类型（int、float、str、list、dict）",
                    "条件语句：if-elif-else结构",
                    "循环语句：for和while循环",
                    "函数定义：def关键字与参数传递"
                ],
                "note": "Python是深度学习领域最常用的编程语言。它语法简洁、易于学习，拥有丰富的科学计算库。掌握Python基础语法是进行深度学习开发的第一步。"
            },
            {
                "title": "列表与字典操作",
                "bullets": [
                    "列表切片：list[start:end:step]",
                    "列表推导式：[x**2 for x in range(10)]",
                    "字典的键值对操作",
                    "常用方法：append、extend、pop、get等"
                ],
                "note": "列表和字典是Python中最常用的数据结构。列表用于存储有序的元素集合，字典用于存储键值对。熟练使用这两种数据结构对于处理深度学习中的数据非常重要。"
            },
            {
                "title": "面向对象编程",
                "bullets": [
                    "类的定义：class关键字",
                    "构造方法：__init__",
                    "实例属性与类属性",
                    "继承与多态"
                ],
                "note": "面向对象编程是构建复杂程序的重要范式。在TensorFlow中，模型通常以类的形式定义。理解类、对象、继承等概念，能帮助我们更好地设计和组织深度学习代码。"
            }
        ],
        "summary": [
            "掌握Python基础语法和数据结构",
            "理解列表推导式等高级特性",
            "能够使用面向对象方式组织代码"
        ]
    },
    {
        "id": "04",
        "title": "科学计算与可视化",
        "subtitle": "第四章 · 神经网络与深度学习",
        "slides": [
            {
                "title": "NumPy数组操作",
                "bullets": [
                    "NumPy是Python科学计算的核心库",
                    "ndarray：高效的多维数组对象",
                    "数组创建：array、zeros、ones、arange",
                    "数组运算：支持广播机制的向量化操作"
                ],
                "note": "NumPy是Python中进行科学计算的基础库。它提供了高效的多维数组对象和丰富的数学函数。在深度学习中，数据通常以NumPy数组的形式进行处理和传递。"
            },
            {
                "title": "NumPy核心操作",
                "bullets": [
                    "数组形状变换：reshape、transpose、flatten",
                    "数组索引与切片：支持多维索引",
                    "数学运算：dot、matmul、sum、mean",
                    "线性代数：矩阵乘法、求逆、特征值分解"
                ],
                "note": "掌握NumPy的核心操作对于深度学习至关重要。reshape用于改变数组形状，矩阵乘法用于神经网络的前向传播，这些操作是构建神经网络的基础。"
            },
            {
                "title": "Matplotlib可视化",
                "bullets": [
                    "plt.plot()：绑制折线图",
                    "plt.scatter()：绘制散点图",
                    "plt.imshow()：显示图像数据",
                    "子图与布局：subplot、figure"
                ],
                "note": "Matplotlib是Python中最常用的可视化库。在深度学习中，我们经常需要可视化训练过程、损失曲线、预测结果等。掌握Matplotlib能帮助我们更好地理解和展示模型性能。"
            }
        ],
        "summary": [
            "熟练使用NumPy进行数组操作",
            "掌握矩阵运算和线性代数基础",
            "能够使用Matplotlib进行数据可视化"
        ]
    },
    {
        "id": "05",
        "title": "TensorFlow基础",
        "subtitle": "第五章 · 神经网络与深度学习",
        "slides": [
            {
                "title": "张量（Tensor）概念",
                "bullets": [
                    "张量是TensorFlow的核心数据结构",
                    "标量（0维）、向量（1维）、矩阵（2维）、高维张量",
                    "创建张量：tf.constant、tf.Variable",
                    "张量属性：shape、dtype、device"
                ],
                "note": "张量是TensorFlow中的基本数据单元，可以理解为多维数组的推广。标量是0维张量，向量是1维张量，矩阵是2维张量。深度学习中的所有数据都以张量形式表示。"
            },
            {
                "title": "张量操作",
                "bullets": [
                    "数学运算：tf.add、tf.multiply、tf.matmul",
                    "形状操作：tf.reshape、tf.transpose",
                    "聚合操作：tf.reduce_sum、tf.reduce_mean",
                    "索引切片：与NumPy语法兼容"
                ],
                "note": "TensorFlow提供了丰富的张量操作函数。这些操作与NumPy非常相似，但可以在GPU上高效执行。掌握这些基本操作是构建神经网络的前提。"
            },
            {
                "title": "自动微分（AutoGrad）",
                "bullets": [
                    "GradientTape：记录计算过程",
                    "自动计算梯度：tape.gradient()",
                    "支持高阶导数计算",
                    "是反向传播算法的核心实现"
                ],
                "note": "自动微分是深度学习框架的核心功能。TensorFlow的GradientTape可以自动记录计算过程并计算梯度，这使得我们无需手动推导复杂的导数公式，大大简化了神经网络的训练过程。"
            }
        ],
        "summary": [
            "理解张量的概念和操作方法",
            "掌握TensorFlow基本运算",
            "理解自动微分的原理和使用"
        ]
    },
    {
        "id": "06",
        "title": "监督学习基础",
        "subtitle": "第六章 · 神经网络与深度学习",
        "slides": [
            {
                "title": "监督学习概述",
                "bullets": [
                    "监督学习：从标注数据中学习映射关系",
                    "训练数据：输入特征X和目标标签Y",
                    "回归问题：预测连续值（如房价预测）",
                    "分类问题：预测离散类别（如图像分类）"
                ],
                "note": "监督学习是机器学习中最常见的范式。它通过学习输入和输出之间的映射关系来进行预测。根据输出类型的不同，可以分为回归问题和分类问题。"
            },
            {
                "title": "线性回归模型",
                "bullets": [
                    "模型形式：y = wx + b",
                    "损失函数：均方误差MSE",
                    "优化目标：最小化损失函数",
                    "梯度下降：迭代更新参数"
                ],
                "note": "线性回归是最简单的监督学习模型。它假设输入和输出之间存在线性关系。通过最小化预测值和真实值之间的均方误差，我们可以学习到最优的权重参数。"
            },
            {
                "title": "梯度下降优化",
                "bullets": [
                    "梯度：损失函数对参数的偏导数",
                    "学习率：控制参数更新步长",
                    "批量梯度下降、随机梯度下降、小批量梯度下降",
                    "收敛条件：损失不再显著下降"
                ],
                "note": "梯度下降是训练神经网络的核心算法。它通过计算损失函数对参数的梯度，沿着梯度的反方向更新参数，逐步找到损失函数的最小值点。"
            }
        ],
        "summary": [
            "理解监督学习的基本概念",
            "掌握线性回归模型的原理",
            "理解梯度下降优化算法"
        ]
    },
    {
        "id": "07",
        "title": "人工神经网络",
        "subtitle": "第七章 · 神经网络与深度学习",
        "slides": [
            {
                "title": "神经元模型",
                "bullets": [
                    "生物神经元的数学抽象",
                    "输入加权求和：z = Σ(wi·xi) + b",
                    "激活函数：引入非线性变换",
                    "常用激活函数：Sigmoid、ReLU、Tanh"
                ],
                "note": "人工神经网络的基本单元是神经元。它模拟生物神经元的工作方式，接收多个输入信号，进行加权求和后通过激活函数产生输出。激活函数引入非线性，使网络能够学习复杂的模式。"
            },
            {
                "title": "多层感知机（MLP）",
                "bullets": [
                    "输入层：接收原始特征",
                    "隐藏层：进行特征变换和提取",
                    "输出层：产生最终预测结果",
                    "全连接结构：每层神经元与下一层全部连接"
                ],
                "note": "多层感知机是最基本的神经网络结构。它由输入层、隐藏层和输出层组成。通过堆叠多个隐藏层，网络可以学习越来越抽象的特征表示。"
            },
            {
                "title": "使用Keras构建MLP",
                "bullets": [
                    "Sequential模型：层的线性堆叠",
                    "Dense层：全连接层",
                    "compile()：配置优化器和损失函数",
                    "fit()：训练模型"
                ],
                "note": "Keras是TensorFlow的高级API，使用它可以非常方便地构建神经网络。只需几行代码，就可以定义、编译和训练一个完整的神经网络模型。"
            }
        ],
        "summary": [
            "理解神经元和激活函数的作用",
            "掌握多层感知机的结构",
            "能够使用Keras构建和训练神经网络"
        ]
    },
    {
        "id": "08",
        "title": "深度学习基础",
        "subtitle": "第八章 · 神经网络与深度学习",
        "slides": [
            {
                "title": "深度学习的优势",
                "bullets": [
                    "自动特征学习：无需手工设计特征",
                    "层次化表示：从低级到高级特征",
                    "端到端学习：直接从原始数据到结果",
                    "大数据时代的最佳选择"
                ],
                "note": "深度学习的最大优势是能够自动从数据中学习特征。传统机器学习需要人工设计特征，而深度学习可以自动发现数据中的模式，这使得它在图像、语音等领域取得了突破性进展。"
            },
            {
                "title": "反向传播算法",
                "bullets": [
                    "前向传播：计算网络输出和损失",
                    "反向传播：从输出层向输入层传递误差",
                    "链式法则：计算每层参数的梯度",
                    "参数更新：使用梯度下降优化"
                ],
                "note": "反向传播是训练神经网络的核心算法。它利用链式法则，从输出层开始逐层向后传播误差，计算每个参数的梯度。这使得我们可以高效地训练拥有数百万参数的深度网络。"
            },
            {
                "title": "常用优化器",
                "bullets": [
                    "SGD：随机梯度下降",
                    "Momentum：带动量的梯度下降",
                    "Adam：自适应学习率优化器",
                    "学习率调度：动态调整学习率"
                ],
                "note": "选择合适的优化器对于训练效果至关重要。Adam是目前最常用的优化器，它结合了动量和自适应学习率的优点，在大多数任务上都能取得很好的效果。"
            },
            {
                "title": "防止过拟合",
                "bullets": [
                    "正则化：L1和L2正则化",
                    "Dropout：随机丢弃神经元",
                    "早停法：监控验证集性能",
                    "数据增强：扩充训练数据"
                ],
                "note": "过拟合是深度学习中的常见问题。当模型在训练集上表现很好但在测试集上表现差时，就发生了过拟合。我们可以通过正则化、Dropout等技术来缓解这个问题。"
            }
        ],
        "summary": [
            "理解深度学习的核心优势",
            "掌握反向传播算法原理",
            "了解常用的优化技术和正则化方法"
        ]
    },
    {
        "id": "09",
        "title": "卷积神经网络",
        "subtitle": "第九章 · 神经网络与深度学习",
        "slides": [
            {
                "title": "卷积操作原理",
                "bullets": [
                    "卷积核：用于提取局部特征的滤波器",
                    "滑动窗口：在输入图像上滑动计算",
                    "特征图：卷积操作的输出",
                    "参数共享：大幅减少参数数量"
                ],
                "note": "卷积操作是CNN的核心。卷积核就像一个小窗口，在图像上滑动并计算局部区域的特征。通过参数共享，CNN可以用很少的参数处理大尺寸图像。"
            },
            {
                "title": "CNN基本组件",
                "bullets": [
                    "卷积层（Conv2D）：提取空间特征",
                    "池化层（MaxPool）：降低特征图尺寸",
                    "全连接层（Dense）：分类决策",
                    "激活函数（ReLU）：引入非线性"
                ],
                "note": "一个典型的CNN由卷积层、池化层和全连接层组成。卷积层提取特征，池化层降低维度，最后全连接层进行分类。这种结构在图像识别任务中非常有效。"
            },
            {
                "title": "经典CNN架构",
                "bullets": [
                    "LeNet：最早的CNN，用于手写数字识别",
                    "AlexNet：2012年ImageNet冠军，开启深度学习时代",
                    "VGG：使用小卷积核的深层网络",
                    "ResNet：残差连接，实现超深网络训练"
                ],
                "note": "了解经典CNN架构对于理解深度学习发展历程很有帮助。从简单的LeNet到深达152层的ResNet，每一次架构创新都推动了图像识别性能的巨大提升。"
            },
            {
                "title": "使用Keras构建CNN",
                "bullets": [
                    "Conv2D：二维卷积层",
                    "MaxPooling2D：最大池化层",
                    "Flatten：展平层，连接卷积和全连接",
                    "实战：MNIST手写数字识别"
                ],
                "note": "使用Keras构建CNN非常简单。我们只需要按照卷积-池化-全连接的顺序堆叠层，然后编译和训练模型。接下来我们将用这个架构完成手写数字识别任务。"
            }
        ],
        "summary": [
            "理解卷积操作的原理和优势",
            "掌握CNN的基本组件和结构",
            "了解经典CNN架构的发展历程"
        ]
    },
    {
        "id": "10",
        "title": "目标检测基础",
        "subtitle": "第十章 · 神经网络与深度学习",
        "slides": [
            {
                "title": "目标检测任务",
                "bullets": [
                    "定位：确定目标在图像中的位置（边界框）",
                    "分类：识别目标的类别",
                    "多目标检测：同时检测多个不同类别的目标",
                    "应用场景：自动驾驶、安防监控、医学影像"
                ],
                "note": "目标检测是计算机视觉中的核心任务之一。与图像分类不同，目标检测不仅要识别图像中有什么物体，还要确定这些物体的精确位置。这项技术在自动驾驶、安防等领域有着广泛应用。"
            },
            {
                "title": "目标检测方法演进",
                "bullets": [
                    "传统方法：滑动窗口 + 手工特征",
                    "两阶段方法：R-CNN系列（先候选区域后分类）",
                    "单阶段方法：YOLO、SSD（端到端检测）",
                    "发展趋势：更快、更准、更轻量"
                ],
                "note": "目标检测方法经历了从传统手工特征到深度学习的演变。两阶段方法如R-CNN精度高但速度慢，单阶段方法如YOLO速度快适合实时应用。目前的研究在平衡速度和精度方面不断取得突破。"
            },
            {
                "title": "YOLO算法简介",
                "bullets": [
                    "You Only Look Once：一次性预测所有目标",
                    "将图像划分为网格，每个网格预测边界框",
                    "实时检测：可达到30FPS以上",
                    "版本演进：YOLOv1到YOLOv8"
                ],
                "note": "YOLO是最著名的单阶段目标检测算法。它将整张图像作为输入，一次性输出所有目标的位置和类别。YOLO以其极快的速度著称，非常适合需要实时检测的应用场景。"
            },
            {
                "title": "目标检测评估指标",
                "bullets": [
                    "IoU（交并比）：预测框与真实框的重合度",
                    "Precision和Recall：精确率和召回率",
                    "mAP（平均精度均值）：综合评估指标",
                    "FPS：每秒处理帧数，衡量速度"
                ],
                "note": "评估目标检测模型需要使用专门的指标。IoU衡量定位精度，mAP综合评估检测性能，FPS反映实时性。在实际应用中，我们需要根据具体需求在精度和速度之间做出权衡。"
            }
        ],
        "summary": [
            "理解目标检测任务的定义和应用",
            "了解主流目标检测算法的发展",
            "掌握目标检测的评估指标"
        ]
    }
]


def generate_all_chapter_ppts(output_dir: str):
    """生成所有章节的PPT"""
    Path(output_dir).mkdir(parents=True, exist_ok=True)

    for chapter in CHAPTERS_DATA:
        output_path = Path(output_dir) / f"chapter_{chapter['id']}_{chapter['title']}.pptx"
        create_chapter_ppt(chapter, str(output_path))

    print(f"\n共生成 {len(CHAPTERS_DATA)} 个章节PPT")


def get_all_voice_scripts():
    """获取所有语音脚本"""
    scripts = []
    for chapter in CHAPTERS_DATA:
        chapter_scripts = {
            "id": chapter["id"],
            "title": chapter["title"],
            "slides": []
        }
        for slide in chapter.get("slides", []):
            if "note" in slide:
                chapter_scripts["slides"].append({
                    "title": slide["title"],
                    "script": slide["note"]
                })
        scripts.append(chapter_scripts)
    return scripts


if __name__ == "__main__":
    output_dir = "/Users/jimfu/Work/huixue(yuanban)/ziyuan_data/课程资源/神经网络与深度学习/assets/slides"
    generate_all_chapter_ppts(output_dir)
