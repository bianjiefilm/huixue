"""
针对 backend/app/services/doc_parser/parser.py 的单元测试。

覆盖:
- parse_pdf / parse_docx / parse_pptx 正例(现场生成临时样例文件)
- 4 类异常路径(损坏文件 / 空文件 / 不支持格式 / 文件不存在),全部断言抛出 DocParseError

不依赖外部资源: PDF 用手写最小合法 PDF 字节流构造(pypdf 可容忍非规范 xref 并
正确提取文字); DOCX/PPTX 用 python-docx / python-pptx 现场生成。
"""

from __future__ import annotations

import os

import pytest

from app.services.doc_parser.parser import (
    DocParseError,
    parse_docx,
    parse_document,
    parse_pdf,
    parse_pptx,
)


# ---------------------------------------------------------------------------
# fixtures: 现场生成的临时样例文件
# ---------------------------------------------------------------------------


def _write_minimal_pdf(path: str, text: str = "Hello World Test PDF Content") -> None:
    """手写一个最小合法 PDF(单页,含可提取文字),不依赖 reportlab。"""
    content = f"BT /F1 24 Tf 72 720 Td ({text}) Tj ET".encode()
    pdf = (
        b"%PDF-1.4\n"
        b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
        b"3 0 obj<</Type/Page/Parent 2 0 R/Resources<</Font<</F1 4 0 R>>>>"
        b"/MediaBox[0 0 612 792]/Contents 5 0 R>>endobj\n"
        b"4 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj\n"
        b"5 0 obj<</Length " + str(len(content)).encode() + b">>\n"
        b"stream\n" + content + b"\nendstream\nendobj\n"
        b"xref\n0 6\ntrailer<</Size 6/Root 1 0 R>>\nstartxref\n0\n%%EOF"
    )
    with open(path, "wb") as f:
        f.write(pdf)


@pytest.fixture
def sample_pdf_path(tmp_path):
    path = str(tmp_path / "sample.pdf")
    _write_minimal_pdf(path)
    return path


@pytest.fixture
def sample_docx_path(tmp_path):
    docx = pytest.importorskip("docx")
    path = str(tmp_path / "sample.docx")
    document = docx.Document()
    document.add_heading("第一章 简介", level=1)
    document.add_paragraph("这是正文第一段。")
    document.add_paragraph("这是正文第二段。")
    document.add_heading("第二章 详情", level=1)
    document.add_paragraph("更多正文内容。")
    table = document.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "列A"
    table.rows[0].cells[1].text = "列B"
    table.rows[1].cells[0].text = "值1"
    table.rows[1].cells[1].text = "值2"
    document.save(path)
    return path


@pytest.fixture
def sample_pptx_path(tmp_path):
    pptx_mod = pytest.importorskip("pptx")
    from pptx.util import Inches

    path = str(tmp_path / "sample.pptx")
    prs = pptx_mod.Presentation()
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "第一页标题"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "副标题内容"

    bullet_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_layout)
    slide2.shapes.title.text = "第二页标题"
    body = slide2.placeholders[1]
    tf = body.text_frame
    tf.text = "要点一"
    p = tf.add_paragraph()
    p.text = "要点二"

    prs.save(path)
    return path


# ---------------------------------------------------------------------------
# 正例: parse_pdf
# ---------------------------------------------------------------------------


def test_parse_pdf_success_structure(sample_pdf_path):
    result = parse_pdf(sample_pdf_path)

    assert result["file_type"] == "pdf"
    assert result["pages"] == 1
    assert isinstance(result["document_id"], str) and result["document_id"].startswith("doc_")
    assert isinstance(result["chunks"], list) and len(result["chunks"]) == 1

    chunk = result["chunks"][0]
    assert chunk["chunk_id"] == "chunk_001"
    assert chunk["page"] == 1
    assert "Hello World Test PDF Content" in chunk["text"]


def test_parse_pdf_multi_page():
    import tempfile

    content1 = b"BT /F1 24 Tf 72 720 Td (Page One Content) Tj ET"
    content2 = b"BT /F1 24 Tf 72 720 Td (Page Two Content) Tj ET"
    pdf = (
        b"%PDF-1.4\n"
        b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Kids[3 0 R 6 0 R]/Count 2>>endobj\n"
        b"3 0 obj<</Type/Page/Parent 2 0 R/Resources<</Font<</F1 4 0 R>>>>"
        b"/MediaBox[0 0 612 792]/Contents 5 0 R>>endobj\n"
        b"4 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj\n"
        b"5 0 obj<</Length " + str(len(content1)).encode() + b">>\nstream\n"
        + content1 + b"\nendstream\nendobj\n"
        b"6 0 obj<</Type/Page/Parent 2 0 R/Resources<</Font<</F1 4 0 R>>>>"
        b"/MediaBox[0 0 612 792]/Contents 7 0 R>>endobj\n"
        b"7 0 obj<</Length " + str(len(content2)).encode() + b">>\nstream\n"
        + content2 + b"\nendstream\nendobj\n"
        b"xref\n0 8\ntrailer<</Size 8/Root 1 0 R>>\nstartxref\n0\n%%EOF"
    )
    with tempfile.TemporaryDirectory() as d:
        path = os.path.join(d, "multi.pdf")
        with open(path, "wb") as f:
            f.write(pdf)
        result = parse_pdf(path)

    assert result["pages"] == 2
    assert len(result["chunks"]) == 2
    assert "Page One Content" in result["chunks"][0]["text"]
    assert "Page Two Content" in result["chunks"][1]["text"]


# ---------------------------------------------------------------------------
# 正例: parse_docx
# ---------------------------------------------------------------------------


def test_parse_docx_success_structure(sample_docx_path):
    result = parse_docx(sample_docx_path)

    assert result["file_type"] == "docx"
    assert result["pages"] == 1
    assert isinstance(result["chunks"], list) and len(result["chunks"]) >= 2

    headings = [c["heading"] for c in result["chunks"]]
    assert "第一章 简介" in headings
    assert "第二章 详情" in headings

    # 表格内容附加为单独 chunk
    table_chunks = [c for c in result["chunks"] if c["heading"] == "表格内容"]
    assert len(table_chunks) == 1
    assert "列A" in table_chunks[0]["text"]
    assert "值1" in table_chunks[0]["text"]


def test_parse_docx_chunk_ids_sequential(sample_docx_path):
    result = parse_docx(sample_docx_path)
    ids = [c["chunk_id"] for c in result["chunks"]]
    assert ids == sorted(ids)
    assert ids[0] == "chunk_001"


# ---------------------------------------------------------------------------
# 正例: parse_pptx
# ---------------------------------------------------------------------------


def test_parse_pptx_success_structure(sample_pptx_path):
    result = parse_pptx(sample_pptx_path)

    assert result["file_type"] == "pptx"
    assert result["pages"] == 2
    assert len(result["chunks"]) == 2

    first = result["chunks"][0]
    assert first["page"] == 1
    assert first["heading"] == "第一页标题"

    second = result["chunks"][1]
    assert second["page"] == 2
    assert second["heading"] == "第二页标题"
    assert "要点一" in second["text"]
    assert "要点二" in second["text"]


def test_parse_pptx_title_from_first_slide(sample_pptx_path):
    result = parse_pptx(sample_pptx_path)
    assert result["title"] == "第一页标题"


# ---------------------------------------------------------------------------
# 异常路径: 文件不存在
# ---------------------------------------------------------------------------


@pytest.mark.parametrize("parser_fn", [parse_pdf, parse_docx, parse_pptx, parse_document])
def test_file_not_found_raises_doc_parse_error(parser_fn, tmp_path):
    missing = str(tmp_path / "does_not_exist.pdf")
    with pytest.raises(DocParseError):
        parser_fn(missing)


# ---------------------------------------------------------------------------
# 异常路径: 空文件 (0 字节)
# ---------------------------------------------------------------------------


@pytest.mark.parametrize("ext", [".pdf", ".docx", ".pptx"])
def test_empty_file_raises_doc_parse_error(ext, tmp_path):
    path = tmp_path / f"empty{ext}"
    path.write_bytes(b"")

    parser_map = {".pdf": parse_pdf, ".docx": parse_docx, ".pptx": parse_pptx}
    with pytest.raises(DocParseError) as exc_info:
        parser_map[ext](str(path))
    assert "空" in str(exc_info.value) or "0" in str(exc_info.value)


# ---------------------------------------------------------------------------
# 异常路径: 损坏文件
# ---------------------------------------------------------------------------


def test_corrupted_pdf_raises_doc_parse_error(tmp_path):
    path = tmp_path / "corrupted.pdf"
    path.write_bytes(b"%PDF-1.4\nthis is not a real pdf structure at all, just garbage bytes" * 5)
    with pytest.raises(DocParseError):
        parse_pdf(str(path))


def test_corrupted_docx_raises_doc_parse_error(tmp_path):
    path = tmp_path / "corrupted.docx"
    # docx 是 zip 包，写入非 zip 垃圾字节触发 PackageNotFoundError
    path.write_bytes(b"this is definitely not a valid zip/docx package" * 5)
    with pytest.raises(DocParseError):
        parse_docx(str(path))


def test_corrupted_pptx_raises_doc_parse_error(tmp_path):
    path = tmp_path / "corrupted.pptx"
    path.write_bytes(b"this is definitely not a valid zip/pptx package" * 5)
    with pytest.raises(DocParseError):
        parse_pptx(str(path))


# ---------------------------------------------------------------------------
# 异常路径: 不支持的文件格式 (通过统一入口 parse_document)
# ---------------------------------------------------------------------------


def test_unsupported_extension_raises_doc_parse_error(tmp_path):
    path = tmp_path / "sample.txt"
    path.write_text("some content", encoding="utf-8")
    with pytest.raises(DocParseError) as exc_info:
        parse_document(str(path))
    assert "不支持" in str(exc_info.value)


def test_no_extension_raises_doc_parse_error(tmp_path):
    path = tmp_path / "sample_no_ext"
    path.write_text("some content", encoding="utf-8")
    with pytest.raises(DocParseError):
        parse_document(str(path))


# ---------------------------------------------------------------------------
# parse_document 统一入口的正例分发
# ---------------------------------------------------------------------------


def test_parse_document_dispatches_to_pdf(sample_pdf_path):
    result = parse_document(sample_pdf_path)
    assert result["file_type"] == "pdf"


def test_parse_document_dispatches_to_docx(sample_docx_path):
    result = parse_document(sample_docx_path)
    assert result["file_type"] == "docx"


def test_parse_document_dispatches_to_pptx(sample_pptx_path):
    result = parse_document(sample_pptx_path)
    assert result["file_type"] == "pptx"
