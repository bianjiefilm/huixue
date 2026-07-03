"""文档解析实现 (PDF / DOCX / PPTX)。

对应《慧学AI升级方案-v2.md》第五章 "资料解析":
- PDF: 提取文字、页码，按页切 chunk（表格提取见后续多模态阶段，本版本不做）
- Word: 提取标题层级、正文，按标题切 chunk
- PPT: 提取每页标题、正文，按页切 chunk

统一返回结构：
{
    "document_id": str,
    "title": str,
    "file_type": "pdf" | "docx" | "pptx",
    "pages": int,
    "chunks": [
        {"chunk_id": str, "page": int, "heading": str | None, "text": str},
        ...
    ],
}

已知局限（见 README）：
- 扫描版 PDF（无文字层）无法提取文字，本模块不做 OCR / 多模态兜底，
  会在返回结构里生成空 chunks 或抛出 DocParseError（视文字量而定）。
- PPTX 解析在未安装 python-pptx 的情况下走 zipfile/XML 兜底，无法提取
  备注(speaker notes)之外的复杂 SmartArt / 嵌入图表文字。
"""

from __future__ import annotations

import os
import re
import uuid
import zipfile
from xml.etree import ElementTree as ET

try:
    from pypdf import PdfReader
    from pypdf.errors import PdfReadError
except ImportError:  # pragma: no cover - pypdf 已在 requirements.txt 中
    PdfReader = None
    PdfReadError = Exception

try:
    import docx as _docx
    from docx.opc.exceptions import PackageNotFoundError
except ImportError:  # pragma: no cover - python-docx 已在 requirements.txt 中
    _docx = None
    PackageNotFoundError = Exception

try:
    from pptx import Presentation  # type: ignore

    _HAS_PYTHON_PPTX = True
except ImportError:  # pptx 库未在 requirements.txt 中，走 zipfile/XML 兜底
    Presentation = None
    _HAS_PYTHON_PPTX = False


class DocParseError(Exception):
    """文档解析失败时抛出的统一异常。

    带清楚的 message，说明文件路径、失败原因，不吞异常。
    """

    def __init__(self, message: str, *, file_path: str | None = None):
        self.file_path = file_path
        full_message = f"{message} (file_path={file_path})" if file_path else message
        super().__init__(full_message)


_SUPPORTED_EXT = {".pdf", ".docx", ".pptx"}

# python-docx heading 样式名匹配，如 "Heading 1" / "标题 1"
_HEADING_STYLE_RE = re.compile(r"^(Heading|标题)\s*([0-9１-９]+)?$", re.IGNORECASE)

# pptx XML namespace
_PPTX_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}


def _new_document_id() -> str:
    return f"doc_{uuid.uuid4().hex[:12]}"


def _new_chunk_id(index: int) -> str:
    return f"chunk_{index:03d}"


def _validate_file(file_path: str) -> None:
    if not file_path or not isinstance(file_path, str):
        raise DocParseError("file_path 不能为空", file_path=file_path)
    if not os.path.isfile(file_path):
        raise DocParseError(f"文件不存在: {file_path}", file_path=file_path)
    if os.path.getsize(file_path) == 0:
        raise DocParseError("文件为空 (0 字节)，无法解析", file_path=file_path)


def _title_from_filename(file_path: str) -> str:
    base = os.path.basename(file_path)
    name, _ext = os.path.splitext(base)
    return name or base


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def parse_pdf(file_path: str) -> dict:
    """解析 PDF 文件，按页切 chunk。

    使用 pypdf（requirements.txt 已有）提取每页文字。
    扫描版 PDF（无文字层）会导致某些页 text 为空字符串，不会抛异常，
    但如果全文档均无法提取到任何文字，会抛 DocParseError 提示疑似扫描版。
    """
    _validate_file(file_path)
    if PdfReader is None:
        raise DocParseError("pypdf 未安装，无法解析 PDF", file_path=file_path)

    try:
        reader = PdfReader(file_path)
    except PdfReadError as exc:
        raise DocParseError(f"PDF 文件损坏或格式非法: {exc}", file_path=file_path) from exc
    except Exception as exc:  # noqa: BLE001 - 需要将底层异常转成统一异常
        raise DocParseError(f"PDF 打开失败: {exc}", file_path=file_path) from exc

    if reader.is_encrypted:
        try:
            # 尝试空密码解密（部分 PDF 仅做了权限限制）
            reader.decrypt("")
        except Exception as exc:  # noqa: BLE001
            raise DocParseError(
                f"PDF 已加密，无法解析: {exc}", file_path=file_path
            ) from exc

    num_pages = len(reader.pages)
    if num_pages == 0:
        raise DocParseError("PDF 页数为 0，无法解析", file_path=file_path)

    document_id = _new_document_id()
    chunks: list[dict] = []
    total_text_len = 0
    chunk_index = 1

    for page_no, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception as exc:  # noqa: BLE001 - 单页提取失败不应中断整份文档
            text = ""
        text = text.strip()
        total_text_len += len(text)

        # 取首行作为该页 heading（若存在）
        heading = None
        if text:
            first_line = text.splitlines()[0].strip()
            if 0 < len(first_line) <= 60:
                heading = first_line

        chunks.append(
            {
                "chunk_id": _new_chunk_id(chunk_index),
                "page": page_no,
                "heading": heading,
                "text": text,
            }
        )
        chunk_index += 1

    if total_text_len == 0:
        raise DocParseError(
            "PDF 全文档未提取到任何文字，疑似扫描版 PDF，需走多模态识别（本模块不支持）",
            file_path=file_path,
        )

    title = _pdf_title(reader, file_path)

    return {
        "document_id": document_id,
        "title": title,
        "file_type": "pdf",
        "pages": num_pages,
        "chunks": chunks,
    }


def _pdf_title(reader, file_path: str) -> str:
    try:
        meta = reader.metadata
        if meta and getattr(meta, "title", None):
            return str(meta.title).strip() or _title_from_filename(file_path)
    except Exception:  # noqa: BLE001
        pass
    return _title_from_filename(file_path)


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def parse_docx(file_path: str) -> dict:
    """解析 DOCX 文件，按标题层级切 chunk。

    使用 python-docx（requirements.txt 已有）。DOCX 没有原生页码概念
    (分页由 Word 渲染时动态计算)，本函数用"段落顺序"近似估计页码：
    page 字段固定为 1，因为无法在不渲染的情况下得到真实页码；
    heading 取最近一个 Heading 样式段落的文本。
    """
    _validate_file(file_path)
    if _docx is None:
        raise DocParseError("python-docx 未安装，无法解析 DOCX", file_path=file_path)

    try:
        document = _docx.Document(file_path)
    except PackageNotFoundError as exc:
        raise DocParseError(
            f"DOCX 文件损坏或不是合法的 docx 包: {exc}", file_path=file_path
        ) from exc
    except Exception as exc:  # noqa: BLE001
        raise DocParseError(f"DOCX 打开失败: {exc}", file_path=file_path) from exc

    document_id = _new_document_id()
    chunks: list[dict] = []
    chunk_index = 1
    current_heading: str | None = None
    current_text_parts: list[str] = []
    any_content = False

    def _flush():
        nonlocal chunk_index
        text = "\n".join(p for p in current_text_parts if p.strip())
        if not text.strip():
            return
        chunks.append(
            {
                "chunk_id": _new_chunk_id(chunk_index),
                "page": 1,
                "heading": current_heading,
                "text": text.strip(),
            }
        )
        chunk_index += 1

    for para in document.paragraphs:
        style_name = (para.style.name if para.style else "") or ""
        text = para.text.strip()
        if not text:
            continue
        any_content = True
        if _HEADING_STYLE_RE.match(style_name):
            # 遇到新标题，先落盘上一个 chunk
            _flush()
            current_heading = text
            current_text_parts = []
        else:
            current_text_parts.append(text)

    _flush()

    # 补充表格文字（不单独切 chunk，附加到最后一个 chunk，或新建一个 chunk）
    table_texts = []
    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                table_texts.append(row_text)
                any_content = True
    if table_texts:
        chunks.append(
            {
                "chunk_id": _new_chunk_id(chunk_index),
                "page": 1,
                "heading": "表格内容",
                "text": "\n".join(table_texts),
            }
        )
        chunk_index += 1

    if not any_content:
        raise DocParseError("DOCX 未提取到任何文字内容（正文与表格均为空）", file_path=file_path)

    title = _docx_title(document, file_path)

    return {
        "document_id": document_id,
        "title": title,
        "file_type": "docx",
        "pages": 1,
        "chunks": chunks,
    }


def _docx_title(document, file_path: str) -> str:
    try:
        core_title = document.core_properties.title
        if core_title and core_title.strip():
            return core_title.strip()
    except Exception:  # noqa: BLE001
        pass
    # 用第一个非空段落作为标题的兜底
    for para in document.paragraphs:
        if para.text.strip():
            return para.text.strip()[:100]
    return _title_from_filename(file_path)


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def parse_pptx(file_path: str) -> dict:
    """解析 PPTX 文件，按页(幻灯片)提取标题与正文。

    优先使用 python-pptx（若已安装）；否则用 zipfile + XML 解析
    ppt/slides/slideN.xml 兜底（不新增依赖）。
    """
    _validate_file(file_path)

    if _HAS_PYTHON_PPTX:
        return _parse_pptx_with_python_pptx(file_path)
    return _parse_pptx_with_zipfile(file_path)


def _parse_pptx_with_python_pptx(file_path: str) -> dict:
    try:
        prs = Presentation(file_path)
    except Exception as exc:  # noqa: BLE001
        raise DocParseError(f"PPTX 打开失败: {exc}", file_path=file_path) from exc

    num_slides = len(prs.slides._sldIdLst)  # noqa: SLF001 - python-pptx 无公开的 slide count API
    if num_slides == 0:
        raise DocParseError("PPTX 幻灯片数为 0，无法解析", file_path=file_path)

    document_id = _new_document_id()
    chunks: list[dict] = []
    chunk_index = 1
    any_content = False
    first_title = None

    for page_no, slide in enumerate(prs.slides, start=1):
        heading = None
        body_parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            shape_text = "\n".join(
                p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()
            )
            if not shape_text:
                continue
            is_title_shape = getattr(shape, "is_placeholder", False) and shape.placeholder_format is not None and shape.placeholder_format.type is not None and shape.placeholder_format.idx == 0
            if heading is None and (is_title_shape or shape == getattr(slide.shapes, "title", None)):
                heading = shape_text.splitlines()[0].strip()
            else:
                body_parts.append(shape_text)

        # 备注
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                body_parts.append(f"[备注] {notes_text}")

        text = "\n".join(body_parts).strip()
        if heading or text:
            any_content = True
        if first_title is None and heading:
            first_title = heading

        chunks.append(
            {
                "chunk_id": _new_chunk_id(chunk_index),
                "page": page_no,
                "heading": heading,
                "text": text,
            }
        )
        chunk_index += 1

    if not any_content:
        raise DocParseError("PPTX 未提取到任何标题或正文内容", file_path=file_path)

    title = first_title or _title_from_filename(file_path)

    return {
        "document_id": document_id,
        "title": title,
        "file_type": "pptx",
        "pages": num_slides,
        "chunks": chunks,
    }


def _parse_pptx_with_zipfile(file_path: str) -> dict:
    """不依赖 python-pptx 的兜底实现：直接读 pptx zip 包内的 slide XML。"""
    try:
        zf = zipfile.ZipFile(file_path)
    except zipfile.BadZipFile as exc:
        raise DocParseError(f"PPTX 文件损坏，不是合法的 zip 包: {exc}", file_path=file_path) from exc
    except Exception as exc:  # noqa: BLE001
        raise DocParseError(f"PPTX 打开失败: {exc}", file_path=file_path) from exc

    with zf:
        slide_names = sorted(
            (n for n in zf.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", n)),
            key=lambda n: int(re.search(r"slide(\d+)\.xml", n).group(1)),
        )
        if not slide_names:
            raise DocParseError("PPTX 内未找到任何幻灯片 XML，无法解析", file_path=file_path)

        document_id = _new_document_id()
        chunks: list[dict] = []
        chunk_index = 1
        any_content = False
        first_title = None

        for page_no, slide_name in enumerate(slide_names, start=1):
            try:
                xml_bytes = zf.read(slide_name)
                root = ET.fromstring(xml_bytes)
            except ET.ParseError as exc:
                raise DocParseError(
                    f"PPTX 幻灯片 XML 解析失败 ({slide_name}): {exc}", file_path=file_path
                ) from exc

            heading, body_lines = _extract_slide_text(root)
            text = "\n".join(body_lines).strip()
            if heading or text:
                any_content = True
            if first_title is None and heading:
                first_title = heading

            chunks.append(
                {
                    "chunk_id": _new_chunk_id(chunk_index),
                    "page": page_no,
                    "heading": heading,
                    "text": text,
                }
            )
            chunk_index += 1

        if not any_content:
            raise DocParseError("PPTX 未提取到任何标题或正文内容", file_path=file_path)

        title = first_title or _title_from_filename(file_path)

        return {
            "document_id": document_id,
            "title": title,
            "file_type": "pptx",
            "pages": len(slide_names),
            "chunks": chunks,
        }


def _extract_slide_text(root: ET.Element) -> tuple[str | None, list[str]]:
    """从 slideN.xml 的 ElementTree 提取标题与正文行。

    slide XML 结构大致为:
    <p:sld><p:cSld><p:spTree><p:sp>
        <p:nvSpPr><p:nvPr><p:ph type="title"/></p:nvPr></p:nvSpPr>
        <p:txBody><a:p><a:r><a:t>文字</a:t></a:r></a:p></p:txBody>
    </p:sp>...
    """
    heading: str | None = None
    body_lines: list[str] = []

    for sp in root.iter("{http://schemas.openxmlformats.org/presentationml/2006/main}sp"):
        is_title = False
        ph = sp.find(".//p:nvSpPr/p:nvPr/p:ph", _PPTX_NS)
        if ph is not None and ph.get("type") in ("title", "ctrTitle"):
            is_title = True

        paragraphs = sp.findall(".//a:txBody/a:p", _PPTX_NS)
        shape_lines = []
        for p in paragraphs:
            runs_text = "".join(
                (t.text or "") for t in p.findall(".//a:t", _PPTX_NS)
            ).strip()
            if runs_text:
                shape_lines.append(runs_text)

        if not shape_lines:
            continue

        if is_title and heading is None:
            heading = shape_lines[0]
            if len(shape_lines) > 1:
                body_lines.extend(shape_lines[1:])
        else:
            body_lines.extend(shape_lines)

    return heading, body_lines


# ---------------------------------------------------------------------------
# 统一入口
# ---------------------------------------------------------------------------

def parse_document(file_path: str) -> dict:
    """按文件扩展名自动分发到对应解析函数。"""
    _validate_file(file_path)
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return parse_pdf(file_path)
    if ext == ".docx":
        return parse_docx(file_path)
    if ext == ".pptx":
        return parse_pptx(file_path)

    raise DocParseError(
        f"不支持的文件格式: {ext or '(无扩展名)'}，仅支持 {sorted(_SUPPORTED_EXT)}",
        file_path=file_path,
    )
